"""Targeted edits to session_deck.pptx for the Q3 FY26 swap and Form 2 reconciliation.

Run this against the live deck. Idempotent at the string level (skips silently if a
string is already updated). Preserves all run-level formatting (font, size, color, bold)
because we only mutate run.text in place.

Slide 4: Zone A one-liner, Zone B header, 5 stat cards swapped to Q3 FY26 actuals.
Slide 5: Surface strip "Mentimeter" -> "MS Forms"; Pair + Post column mirrors Form 2 fields.
Slide 7: Step 5 + deliverable line mirror Form 2 fields.
Slide 8: Step 2 body mirrors Form 2 fields.
"""

from pptx import Presentation

DECK = "/sessions/peaceful-gallant-ramanujan/mnt/earnings-demo/workshop/session_deck.pptx"


def replace_run_text(slide, find, replace):
    """Replace the exact text of any run that matches `find`. Returns count."""
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == find:
                    run.text = replace
                    count += 1
    return count


def replace_substring(slide, find, replace):
    """Replace `find` substring anywhere it appears in a run. Returns count."""
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    count += 1
    return count


def update_stat_card(slide, label_text, new_value, new_sub):
    """Find the stat card whose first paragraph is `label_text`, update value (paragraph[1])
    and sub (paragraph[2])."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = shape.text_frame.paragraphs
        if len(paras) < 3:
            continue
        first_text = "".join(r.text for r in paras[0].runs).strip()
        if first_text != label_text:
            continue
        # Found the card. Update paragraph[1] and paragraph[2].
        # Each typically has one run. We update the first run and clear any extras.
        for para_idx, new_text in [(1, new_value), (2, new_sub)]:
            para = paras[para_idx]
            if not para.runs:
                continue
            para.runs[0].text = new_text
            # Clear additional runs in this paragraph to avoid stale fragments
            for extra in para.runs[1:]:
                extra.text = ""
        return True
    return False


def main():
    p = Presentation(DECK)

    # Slide 4 (index 3): Problem Setup
    slide4 = p.slides[3]
    log = []

    # Zone A one-liner: swap "Q2 FY26 in February" → "Q3 FY26 on June 2"
    n = replace_substring(slide4, "Q2 FY26 in February", "Q3 FY26 on June 2")
    log.append(f"Slide 4 Zone A one-liner swap: {n} run(s)")

    # Zone B header: "Q2 FY26 print (February 2026)  ·  refresh on June 3 with Q3 FY26 figures"
    n = replace_substring(
        slide4,
        "Q2 FY26 print (February 2026)  ·  refresh on June 3 with Q3 FY26 figures",
        "Q3 FY26 print (May 2026)  ·  reported June 2",
    )
    log.append(f"Slide 4 Zone B header swap: {n} run(s)")

    # Stat cards
    cards = [
        ("Revenue", "$3.00B", "+31.1% YoY (+2.1% beat)"),
        ("NGS ARR", "$8.13B", "+60% reported / +28% organic"),
        ("Non-GAAP EPS", "$0.85", "vs $0.80 consensus (+6.2%)"),
        ("Stock reaction", "-4.4%", "next-day gap after print"),
        ("The twist", "Reported vs organic", "+60% (M&A-boosted) vs +28% organic"),
    ]
    for label, value, sub in cards:
        ok = update_stat_card(slide4, label, value, sub)
        log.append(f"Slide 4 stat card '{label}': {'updated' if ok else 'NOT FOUND'}")

    # Zone C task block: update to mirror Form 2 fields
    n = replace_substring(
        slide4,
        "Plus one sentence on your biggest conviction. Plus one sentence on your biggest uncertainty.",
        "Plus your confidence (1-5), primary reason, and biggest risk.",
    )
    log.append(f"Slide 4 task block reconcile: {n} run(s)")

    # Slide 5 (index 4): Three Beats
    slide5 = p.slides[4]
    n = replace_run_text(slide5, "Mentimeter", "MS Forms")
    log.append(f"Slide 5 surface strip Mentimeter->MS Forms: {n} run(s)")
    n = replace_run_text(slide5, "One sentence biggest conviction.", "Confidence (1-5).")
    log.append(f"Slide 5 Pair+Post conviction line: {n} run(s)")
    n = replace_run_text(slide5, "One sentence biggest uncertainty.", "Primary reason. Biggest risk.")
    log.append(f"Slide 5 Pair+Post uncertainty line: {n} run(s)")

    # Slide 7 (index 6): Beat 2 standing
    slide7 = p.slides[6]
    n = replace_substring(
        slide7,
        "Land it. Buy / Hold / Sell with one-sentence conviction and uncertainty.",
        "Land it. Buy / Hold / Sell with confidence, primary reason, and biggest risk.",
    )
    log.append(f"Slide 7 step 5: {n} run(s)")
    n = replace_substring(
        slide7,
        "Deliverable: Buy / Hold / Sell + one sentence conviction + one sentence uncertainty.",
        "Deliverable: Buy / Hold / Sell + confidence + primary reason + biggest risk.",
    )
    log.append(f"Slide 7 deliverable line: {n} run(s)")

    # Slide 8 (index 7): Pair + Post standing
    slide8 = p.slides[7]
    n = replace_substring(
        slide8,
        "Buy / Hold / Sell. One sentence conviction. One sentence uncertainty.",
        "Buy / Hold / Sell. Confidence 1-5. Primary reason. Biggest risk.",
    )
    log.append(f"Slide 8 step 2 body: {n} run(s)")

    p.save(DECK)
    print("Saved.")
    for line in log:
        print(" ", line)


if __name__ == "__main__":
    main()
