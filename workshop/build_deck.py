"""Build the workshop session deck on the Aileron template.

Reads:  session_deck_structured.pptx (template with slides duplicated in order)
Writes: session_deck.pptx (final deck)

Slide map (final order):
  1  Title                          (kept from template slide 1, edited in XML)
  2  Settling poll                  (was AG slide 6: Bullets + Image)
  3  5 components                   (was AG slide 5: Complex - rebuilt as table)
  4  The call and your question     (was AG slide 5 copy: Complex - rebuilt as 3 zones)
  5  Three Beats                    (was AG slide 4: Framework Matrix)
  6  Beat 1 standing                (was AG slide 7: Icon + Text Rows)
  7  Beat 2 standing                (was AG slide 5 copy: Complex - rebuilt as 3 zones)
  8  Pair + Post standing           (was AG slide 6 copy: Bullets + Image)
  9  Debrief                        (was AG slide 4 copy: Framework Matrix)
  10 Take-home                      (was AG slide 7 copy: Icon + Text Rows)
  11 Contact                        (AG slide 10, static)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# AG brand colors
AG_BLUE = RGBColor(0x60, 0xB5, 0xE5)
AG_DARK = RGBColor(0x2D, 0x20, 0x42)
AG_LIGHT_BLUE = RGBColor(0xB3, 0xDC, 0xF3)
AG_OFFWHITE = RGBColor(0xF2, 0xF2, 0xF2)
AG_BLACK = RGBColor(0x00, 0x00, 0x00)
AG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AG_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

INPUT = "session_deck_structured.pptx"
OUTPUT = "session_deck.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def strip_slide_keep_footer(slide):
    """Remove all non-essential shapes, keeping: footer text, page number, watermark image, slide number placeholder.
    Returns list of shapes to keep so caller knows what is preserved.
    """
    keep_names = ('Footer Placeholder', 'Slide Number Placeholder', 'AI watermark', 'Ai Icon', 'aileron-ai')
    shapes_to_remove = []
    for shape in list(slide.shapes):
        name = shape.name or ''
        # Keep placeholders for footer/page/watermark
        is_keeper = False
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # 14=footer, 13=slide number, 12=date
            if ph_type in (14, 13, 12):
                is_keeper = True
        # Keep AI watermark image (small image bottom-left)
        if shape.shape_type == 13:  # picture
            if shape.top and shape.top > Inches(6.5) and shape.left < Inches(1):
                if shape.width and shape.width < Inches(1):
                    is_keeper = True
        for kn in keep_names:
            if kn.lower() in name.lower():
                is_keeper = True
        if not is_keeper:
            shapes_to_remove.append(shape)
    # Remove
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_blue_header_bar(slide, title_text, subtitle_text=None, height_in=1.0, title_size=28):
    """Add a blue header bar at the top with title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(height_in)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AG_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.7)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = 'Montserrat'
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = AG_WHITE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.name = 'Montserrat'
        r2.font.size = Pt(13)
        r2.font.color.rgb = AG_WHITE
    return bar


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, font='Montserrat', anchor=MSO_ANCHOR.TOP):
    """Add a text box with single paragraph of formatted text."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    else:
        # list of (text, size, bold, color) tuples per paragraph
        for i, item in enumerate(text):
            if i == 0:
                pp = p
            else:
                pp = tf.add_paragraph()
            pp.alignment = align
            if isinstance(item, str):
                r = pp.add_run()
                r.text = item
                r.font.name = font
                r.font.size = Pt(size)
                r.font.bold = bold
                if color is not None:
                    r.font.color.rgb = color
            elif isinstance(item, dict):
                r = pp.add_run()
                r.text = item.get('text', '')
                r.font.name = item.get('font', font)
                r.font.size = Pt(item.get('size', size))
                r.font.bold = item.get('bold', bold)
                if item.get('color') is not None:
                    r.font.color.rgb = item['color']
    return box


def add_rich_text(slide, left, top, width, height, paragraphs, *, anchor=MSO_ANCHOR.TOP):
    """Add a textbox with multiple paragraphs.

    paragraphs: list of dicts: {text, size, bold, color, align, font, space_after}
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = para.get('align', PP_ALIGN.LEFT)
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        if para.get('space_before') is not None:
            p.space_before = Pt(para['space_before'])
        r = p.add_run()
        r.text = para.get('text', '')
        r.font.name = para.get('font', 'Montserrat')
        r.font.size = Pt(para.get('size', 14))
        r.font.bold = para.get('bold', False)
        r.font.italic = para.get('italic', False)
        if para.get('color') is not None:
            r.font.color.rgb = para['color']
    return box


def add_rect(slide, left, top, width, height, fill=None, line=None, line_width=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_width:
            rect.line.width = line_width
    rect.text_frame.text = ''
    return rect


def add_rounded_rect(slide, left, top, width, height, fill=None, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
    return rect


def add_qr_placeholder(slide, left, top, size, label):
    """Add a placeholder square for QR code with label below."""
    # Outer box with light blue fill
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size)
    box.fill.solid()
    box.fill.fore_color.rgb = AG_OFFWHITE
    box.line.color.rgb = AG_DARK
    box.line.width = Pt(1)
    # Add "QR" text inside
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = 'QR'
    r1.font.name = 'Montserrat'
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = AG_DARK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = '[placeholder]'
    r2.font.name = 'Montserrat'
    r2.font.size = Pt(9)
    r2.font.italic = True
    r2.font.color.rgb = AG_DARK
    # Label below
    add_text(
        slide,
        left, top + size + Inches(0.05),
        size, Inches(0.4),
        label,
        size=10, color=AG_DARK_GRAY, align=PP_ALIGN.CENTER
    )


def update_footer_page(slide, page_num):
    """Update the page number text in the footer area."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type == 13:  # slide number
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(page_num)
            r.font.name = 'Montserrat'
            r.font.size = Pt(10)
            r.font.color.rgb = AG_DARK_GRAY


def find_and_clear_page_number_text(slide, new_page_num):
    """For hardcoded page numbers as plain text boxes (not placeholders).

    Look for small text boxes in the bottom-right with a numeric value.
    """
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder:
            continue
        # Check if positioned in footer area
        if not (shape.top and shape.top > Inches(6.7)):
            continue
        text = shape.text_frame.text.strip()
        if text.isdigit() or (text and len(text) <= 3 and any(c.isdigit() for c in text)):
            # Replace with new number
            tf = shape.text_frame
            # Preserve formatting if possible: get first run's formatting
            font_name = 'Montserrat'
            font_size = 10
            font_color = AG_DARK_GRAY
            try:
                first_run = tf.paragraphs[0].runs[0]
                font_name = first_run.font.name or font_name
                if first_run.font.size:
                    font_size = first_run.font.size.pt
                if first_run.font.color and first_run.font.color.rgb:
                    font_color = first_run.font.color.rgb
            except (IndexError, AttributeError):
                pass
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = str(new_page_num)
            r.font.name = font_name
            r.font.size = Pt(font_size)
            r.font.color.rgb = font_color


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide_2_settling_poll(slide):
    """Slide 2: Settling poll. 3 questions + QR placeholder."""
    strip_slide_keep_footer(slide)
    add_blue_header_bar(slide, 'While you settle in', 'Three questions. Tap to respond on your phone.')

    # Three questions on the left
    questions = [
        ('1', 'How would you describe your current use of AI in your professional work?'),
        ('2', 'Recent market reactions to domain-specific AI launches (Claude for Legal, AI for financial services). How do you read them?'),
        ('3', 'Will AI be net positive for investment analysis specifically?'),
    ]
    top = Inches(1.4)
    for num, q in questions:
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), top, Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = AG_LIGHT_BLUE
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.margin_top = Inches(0)
        ctf.margin_bottom = Inches(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.name = 'Montserrat'
        cr.font.size = Pt(20)
        cr.font.bold = True
        cr.font.color.rgb = AG_DARK
        # Question text
        add_text(
            slide,
            Inches(1.65), top + Inches(0.02),
            Inches(7.5), Inches(1.4),
            q,
            size=18, bold=True, color=AG_DARK,
            anchor=MSO_ANCHOR.TOP,
        )
        top += Inches(1.5)

    # QR placeholder right side
    add_qr_placeholder(
        slide,
        Inches(10.0), Inches(2.0),
        Inches(2.4),
        'Mentimeter settling poll'
    )

    # Callout bar at bottom
    callout_top = Inches(6.3)
    callout = add_rect(slide, Inches(0.5), callout_top, Inches(12.3), Inches(0.5), fill=AG_LIGHT_BLUE)
    tf = callout.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Anonymous. No name. No login. We will reveal where the room landed in a few minutes.'
    r.font.name = 'Montserrat'
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = AG_DARK


def build_slide_3_five_components(slide):
    """Slide 3: 5 components (Behind the Veil core). Table-style layout."""
    strip_slide_keep_footer(slide)
    add_blue_header_bar(slide, "What's actually happening behind the prompt", 'Key concepts · practitioner fluency', height_in=1.1, title_size=26)

    # Table: 6 rows (header + 5 components), 3 columns
    # Col widths: 3.4" / 6.0" / 3.4"
    col1_w = Inches(3.5)
    col2_w = Inches(5.7)
    col3_w = Inches(3.6)
    col1_x = Inches(0.4)
    col2_x = col1_x + col1_w
    col3_x = col2_x + col2_w

    # Header row
    header_top = Inches(1.35)
    header_h = Inches(0.45)
    headers = [
        (col1_x, col1_w, 'COMPONENT'),
        (col2_x, col2_w, "WHAT'S ACTUALLY HAPPENING"),
        (col3_x, col3_w, 'WHAT PRACTITIONERS NEED'),
    ]
    for x, w, label in headers:
        rect = add_rect(slide, x, header_top, w, header_h, fill=AG_DARK)
        tf = rect.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = label
        r.font.name = 'Montserrat'
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = AG_WHITE

    # 5 component rows
    components = [
        ('01', 'The Interface',
         'ChatGPT, Claude, Copilot, apps built on them. The surface where prompts and outputs exchange. Quality of output is bounded by quality of input.',
         'DESCRIPTION',
         'Define the problem precisely. Articulate intent, scope, success criteria.'),
        ('02', 'The Model (LLM)',
         'Stochastic prediction of the next token from training patterns. Real capabilities, real limitations. Reasoning is pattern-completion, not deduction.',
         'DISCERNMENT',
         'Know what the model can and cannot do. Select the right model. Know when not to use one.'),
        ('03', 'Knowledge & Data Access',
         'How models reach beyond training: RAG, search, documents, databases. The model amplifies whatever you feed it.',
         'DATA LITERACY',
         'Own the quality, access, and governance of your data. Garbage in, garbage out, at enterprise scale.'),
        ('04', 'Agents',
         'Multi-step workflows, tool use, autonomous task completion. Agents take actions: send the email, modify the record, move the money. The risk surface is the action, not the answer.',
         'DELEGATION',
         'Manage agents like junior analysts. Scope the task. Set guardrails. Review the output. Trust is earned.'),
        ('05', 'Guardrails',
         'Safety, accuracy, compliance, organizational policy. Controls that turn capability into accountable capability.',
         'DILIGENCE',
         'Every output is verified. Standard is not perfection but accountability. Human in the loop. Documented. Auditable.'),
    ]
    row_top = header_top + header_h
    row_h = Inches(0.93)
    for num, comp_name, happening, competency, need in components:
        # Col 1: number + component name
        c1 = add_rect(slide, col1_x, row_top, col1_w, row_h, line=AG_LIGHT_BLUE, line_width=Pt(0.5))
        tf1 = c1.text_frame
        tf1.margin_left = Inches(0.18)
        tf1.margin_top = Inches(0.1)
        tf1.margin_right = Inches(0.1)
        tf1.vertical_anchor = MSO_ANCHOR.TOP
        # Number on first paragraph
        p_num = tf1.paragraphs[0]
        p_num.alignment = PP_ALIGN.LEFT
        r_num = p_num.add_run()
        r_num.text = num
        r_num.font.name = 'Montserrat'
        r_num.font.size = Pt(20)
        r_num.font.bold = True
        r_num.font.color.rgb = AG_BLUE
        # Component name
        p_name = tf1.add_paragraph()
        p_name.alignment = PP_ALIGN.LEFT
        p_name.space_before = Pt(2)
        r_name = p_name.add_run()
        r_name.text = comp_name
        r_name.font.name = 'Montserrat'
        r_name.font.size = Pt(13)
        r_name.font.bold = True
        r_name.font.color.rgb = AG_DARK

        # Col 2: What's actually happening
        c2 = add_rect(slide, col2_x, row_top, col2_w, row_h, line=AG_LIGHT_BLUE, line_width=Pt(0.5))
        tf2 = c2.text_frame
        tf2.margin_left = Inches(0.15)
        tf2.margin_top = Inches(0.1)
        tf2.margin_right = Inches(0.15)
        tf2.margin_bottom = Inches(0.05)
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.TOP
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = happening
        r2.font.name = 'Montserrat'
        r2.font.size = Pt(10)
        r2.font.color.rgb = AG_DARK_GRAY

        # Col 3: Competency tag + need
        c3 = add_rect(slide, col3_x, row_top, col3_w, row_h, line=AG_LIGHT_BLUE, line_width=Pt(0.5))
        tf3 = c3.text_frame
        tf3.margin_left = Inches(0.15)
        tf3.margin_top = Inches(0.1)
        tf3.margin_right = Inches(0.15)
        tf3.margin_bottom = Inches(0.05)
        tf3.word_wrap = True
        tf3.vertical_anchor = MSO_ANCHOR.TOP
        p3a = tf3.paragraphs[0]
        p3a.alignment = PP_ALIGN.LEFT
        r3a = p3a.add_run()
        r3a.text = competency
        r3a.font.name = 'Montserrat'
        r3a.font.size = Pt(11)
        r3a.font.bold = True
        r3a.font.color.rgb = AG_BLUE
        p3b = tf3.add_paragraph()
        p3b.alignment = PP_ALIGN.LEFT
        p3b.space_before = Pt(2)
        r3b = p3b.add_run()
        r3b.text = need
        r3b.font.name = 'Montserrat'
        r3b.font.size = Pt(10)
        r3b.font.color.rgb = AG_DARK_GRAY

        row_top += row_h

    # Footer line: 4D + Data Literacy
    add_text(
        slide,
        Inches(0.4), Inches(6.8),
        Inches(12.5), Inches(0.3),
        'The five practitioner competencies extend Anthropic’s 4D AI Fluency Framework (Description, Discernment, Delegation, Diligence) with Data Literacy for the enterprise context.',
        size=10, italic_safe=False, color=AG_DARK_GRAY, align=PP_ALIGN.LEFT, font='Montserrat',
    ) if False else add_rich_text(
        slide,
        Inches(0.4), Inches(6.8),
        Inches(12.5), Inches(0.3),
        [{
            'text': 'The five practitioner competencies extend Anthropic’s 4D AI Fluency Framework (Description, Discernment, Delegation, Diligence) with Data Literacy for the enterprise context.',
            'size': 10, 'italic': True, 'color': AG_DARK_GRAY,
        }]
    )


def build_slide_4_call_and_question(slide):
    """Slide 4: The call and your question. Three zones: one-liner, highlights box, task+QR."""
    strip_slide_keep_footer(slide)
    add_blue_header_bar(slide, 'The call and your question')

    # Zone A: One-liner + why PANW
    add_rich_text(
        slide,
        Inches(0.5), Inches(1.15),
        Inches(12.3), Inches(1.0),
        [
            {'text': 'Palo Alto Networks reported Q2 FY26 in February. Your task: form an investment view.',
             'size': 22, 'bold': True, 'color': AG_DARK, 'space_after': 4},
            {'text': 'AI governance and cybersecurity tracks on the forum agenda. The market just told you what it thinks. The question is what you think.',
             'size': 13, 'italic': True, 'color': AG_DARK_GRAY},
        ]
    )

    # Zone B: 5-slot highlights box (Q2 FY26 actuals)
    box_top = Inches(2.5)
    box_h = Inches(2.6)
    # Container
    container = add_rect(slide, Inches(0.5), box_top, Inches(12.3), box_h, fill=AG_OFFWHITE, line=AG_LIGHT_BLUE, line_width=Pt(1))
    container.text_frame.text = ''

    # Header for the box
    add_rich_text(
        slide,
        Inches(0.7), box_top + Inches(0.1),
        Inches(12.0), Inches(0.4),
        [{
            'text': 'Q2 FY26 print (February 2026)  ·  refresh on June 3 with Q3 FY26 figures',
            'size': 12, 'bold': True, 'color': AG_DARK,
        }]
    )

    # 5 stat slots in a row (or 2 rows if needed)
    stats = [
        ('Revenue', '$2.59B', '+14.9% YoY organic'),
        ('NGS ARR', '$6.33B', '+33% YoY / +28% organic'),
        ('Non-GAAP EPS', '$1.03', 'vs $0.94 consensus (+9.9%)'),
        ('Stock reaction', '-8.5%', 'next-day gap after print'),
        ('The twist', 'Beat then guide-down', 'Q3 EPS guide $0.79 vs Q2 actual $1.03'),
    ]
    n = len(stats)
    available_w = Inches(11.9)
    gap = Inches(0.15)
    cell_w = Emu(int((available_w - gap * (n - 1)) / n))
    cell_top = box_top + Inches(0.6)
    cell_h = Inches(1.9)
    x = Inches(0.7)
    for label, value, sub in stats:
        # Background card
        card = add_rect(slide, x, cell_top, cell_w, cell_h, fill=AG_WHITE, line=AG_BLUE, line_width=Pt(0.75))
        tf = card.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.1)
        tf.word_wrap = True
        # Label
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label
        r1.font.name = 'Montserrat'
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = AG_DARK
        # Value
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = value
        r2.font.name = 'Montserrat'
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = AG_BLUE
        # Sub
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = sub
        r3.font.name = 'Montserrat'
        r3.font.size = Pt(9)
        r3.font.color.rgb = AG_DARK_GRAY
        x = x + cell_w + gap

    # Zone C: Task + QR (bottom)
    add_rich_text(
        slide,
        Inches(0.5), Inches(5.3),
        Inches(9.0), Inches(1.4),
        [
            {'text': 'Your task',
             'size': 12, 'bold': True, 'color': AG_BLUE, 'space_after': 2},
            {'text': 'Form an investment view: Buy / Hold / Sell.',
             'size': 20, 'bold': True, 'color': AG_DARK, 'space_after': 4},
            {'text': 'Plus one sentence on your biggest conviction. Plus one sentence on your biggest uncertainty.',
             'size': 13, 'color': AG_DARK_GRAY},
        ]
    )

    # QR placeholder
    add_qr_placeholder(
        slide,
        Inches(10.0), Inches(5.3),
        Inches(1.4),
        'PANW one-pager'
    )


def build_slide_5_three_beats(slide):
    """Slide 5: Three Beats. 3-column timeline with surface-switching strip."""
    strip_slide_keep_footer(slide)
    add_blue_header_bar(slide, 'How the next 30 minutes will run', 'Three beats. Lead with your framework.')

    # Three columns
    col_top = Inches(1.4)
    col_h = Inches(4.4)
    margin = Inches(0.4)
    gap = Inches(0.25)
    total_w = SLIDE_W - margin * 2
    col_w = Emu(int((total_w - gap * 2) / 3))

    beats = [
        {
            'badge': 'BEAT 1',
            'title': 'Think before you type',
            'time': '4 min  ·  no laptops',
            'icon': 'No laptops',
            'prompt_lead': 'Two questions:',
            'items': [
                'What factors drive your investment view on this company?',
                'What would you want to know?',
            ],
            'note': 'Write it down. You will need it.',
        },
        {
            'badge': 'BEAT 2',
            'title': 'Work the problem with AI',
            'time': '12 min  ·  laptops and phones',
            'icon': 'Pair up',
            'prompt_lead': 'Lead with your framework:',
            'items': [
                '1. Share your framework. Ask what is missing.',
                '2. Decide what data you need. Pull from the one-pager. Ask AI for the rest.',
                '3. Pick your biggest uncertainty. Drill in.',
                '4. Ask AI to argue the strongest case against your view.',
                '5. Land it.',
            ],
            'note': 'You will not finish all five. That is fine.',
        },
        {
            'badge': 'PAIR + POST',
            'title': 'Compare and post',
            'time': '4 min  ·  combined view',
            'icon': 'Anonymous',
            'prompt_lead': 'With a neighbor:',
            'items': [
                'Land one combined view.',
                'Post Buy / Hold / Sell.',
                'One sentence biggest conviction.',
                'One sentence biggest uncertainty.',
            ],
            'note': 'Then read the room. We aggregate live.',
        },
    ]

    x = margin
    for beat in beats:
        # Badge
        badge_h = Inches(0.4)
        badge = add_rect(slide, x, col_top, col_w, badge_h, fill=AG_BLUE)
        btf = badge.text_frame
        btf.margin_left = Inches(0.15)
        btf.margin_top = Inches(0)
        btf.margin_bottom = Inches(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run()
        br.text = beat['badge']
        br.font.name = 'Montserrat'
        br.font.size = Pt(13)
        br.font.bold = True
        br.font.color.rgb = AG_WHITE

        # Body
        body_top = col_top + badge_h + Inches(0.05)
        body_h = col_h - badge_h - Inches(0.05)
        body = add_rect(slide, x, body_top, col_w, body_h, fill=AG_OFFWHITE)

        # Build paragraphs in the body text frame
        tf = body.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # Title
        p_title = tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.LEFT
        r_title = p_title.add_run()
        r_title.text = beat['title']
        r_title.font.name = 'Montserrat'
        r_title.font.size = Pt(16)
        r_title.font.bold = True
        r_title.font.color.rgb = AG_DARK

        # Time
        p_time = tf.add_paragraph()
        p_time.alignment = PP_ALIGN.LEFT
        p_time.space_after = Pt(8)
        r_time = p_time.add_run()
        r_time.text = beat['time']
        r_time.font.name = 'Montserrat'
        r_time.font.size = Pt(11)
        r_time.font.italic = True
        r_time.font.color.rgb = AG_BLUE

        # Prompt lead
        p_lead = tf.add_paragraph()
        p_lead.alignment = PP_ALIGN.LEFT
        p_lead.space_before = Pt(6)
        p_lead.space_after = Pt(4)
        r_lead = p_lead.add_run()
        r_lead.text = beat['prompt_lead']
        r_lead.font.name = 'Montserrat'
        r_lead.font.size = Pt(11)
        r_lead.font.bold = True
        r_lead.font.color.rgb = AG_DARK

        # Items
        for it in beat['items']:
            p_it = tf.add_paragraph()
            p_it.alignment = PP_ALIGN.LEFT
            p_it.space_after = Pt(2)
            r_it = p_it.add_run()
            r_it.text = it
            r_it.font.name = 'Montserrat'
            r_it.font.size = Pt(10)
            r_it.font.color.rgb = AG_DARK_GRAY

        # Note
        p_note = tf.add_paragraph()
        p_note.alignment = PP_ALIGN.LEFT
        p_note.space_before = Pt(10)
        r_note = p_note.add_run()
        r_note.text = beat['note']
        r_note.font.name = 'Montserrat'
        r_note.font.size = Pt(10)
        r_note.font.italic = True
        r_note.font.color.rgb = AG_DARK

        x = x + col_w + gap

    # Bottom: surface-switching strip
    strip_top = Inches(6.0)
    strip_h = Inches(0.4)
    surfaces = [
        ('Beat 1', 'Deck'),
        ('Beat 2', 'Deck'),
        ('Pair + Post', 'Deck'),
        ('Read', 'Mentimeter'),
        ('Demo', 'Cowork + HTML'),
        ('Debrief', 'Deck'),
    ]
    seg_w = Emu(int((SLIDE_W - Inches(0.8)) / len(surfaces)))
    x = Inches(0.4)
    for label, surface in surfaces:
        is_external = surface != 'Deck'
        fill = AG_DARK if is_external else AG_LIGHT_BLUE
        text_color = AG_WHITE if is_external else AG_DARK
        seg = add_rect(slide, x, strip_top, seg_w - Inches(0.04), strip_h, fill=fill)
        tf = seg.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = 'Montserrat'
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = text_color
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = surface
        r2.font.name = 'Montserrat'
        r2.font.size = Pt(7)
        r2.font.color.rgb = text_color
        x = x + seg_w

    # Footer line
    add_rich_text(
        slide,
        Inches(0.4), Inches(6.5),
        Inches(12.5), Inches(0.3),
        [{
            'text': 'There is no right answer. The point is the quality of your reasoning.',
            'size': 11, 'italic': True, 'color': AG_DARK, 'align': PP_ALIGN.CENTER,
        }]
    )


def build_slide_6_beat_1(slide):
    """Slide 6: Beat 1 standing. Text-dominant, two big questions."""
    strip_slide_keep_footer(slide)
    # No blue bar on this archetype - just title in AG blue text
    add_rich_text(
        slide,
        Inches(0.7), Inches(0.5),
        Inches(12.0), Inches(0.7),
        [{
            'text': 'Beat 1  ·  Think before you type',
            'size': 22, 'bold': True, 'color': AG_BLUE,
        }]
    )

    # Time + no-laptops indicator
    add_rich_text(
        slide,
        Inches(0.7), Inches(1.2),
        Inches(12.0), Inches(0.4),
        [{
            'text': '4 minutes  ·  No laptops. No phones. Just you and the page.',
            'size': 14, 'italic': True, 'color': AG_DARK_GRAY,
        }]
    )

    # Two big questions
    questions = [
        'What factors drive your investment view on this company?',
        'What would you want to know?',
    ]

    q_top = Inches(2.0)
    q_h = Inches(1.6)
    for i, q in enumerate(questions):
        # Container with light fill
        box = add_rect(slide, Inches(0.7), q_top, Inches(11.9), q_h, fill=AG_OFFWHITE, line=AG_LIGHT_BLUE, line_width=Pt(1))
        tf = box.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Q number
        p_num = tf.paragraphs[0]
        p_num.alignment = PP_ALIGN.LEFT
        r_num = p_num.add_run()
        r_num.text = f'Q{i+1}'
        r_num.font.name = 'Montserrat'
        r_num.font.size = Pt(14)
        r_num.font.bold = True
        r_num.font.color.rgb = AG_BLUE
        # Question
        p_q = tf.add_paragraph()
        p_q.alignment = PP_ALIGN.LEFT
        p_q.space_before = Pt(4)
        r_q = p_q.add_run()
        r_q.text = q
        r_q.font.name = 'Montserrat'
        r_q.font.size = Pt(28)
        r_q.font.bold = True
        r_q.font.color.rgb = AG_DARK
        q_top += q_h + Inches(0.2)

    # Bottom line
    add_rich_text(
        slide,
        Inches(0.7), Inches(5.7),
        Inches(11.9), Inches(0.5),
        [{
            'text': 'Write it down. You will need it.',
            'size': 16, 'italic': True, 'color': AG_DARK_GRAY, 'align': PP_ALIGN.CENTER,
        }]
    )


def build_slide_7_beat_2(slide):
    """Slide 7: Beat 2 standing. Option C: banner + sequence + footer strip."""
    strip_slide_keep_footer(slide)
    # Top banner: Pairing rule
    banner_h = Inches(0.9)
    banner = add_rect(slide, Inches(0), Inches(0), SLIDE_W, banner_h, fill=AG_DARK)
    tf = banner.text_frame
    tf.margin_left = Inches(0.7)
    tf.margin_right = Inches(0.7)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = 'Beat 2  ·  Work the problem with AI  ·  12 minutes'
    r1.font.name = 'Montserrat'
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = AG_LIGHT_BLUE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = 'Pair up: laptop users raise hands, pair with phone-only. Phone person articulates. Laptop person prompts.'
    r2.font.name = 'Montserrat'
    r2.font.size = Pt(13)
    r2.font.color.rgb = AG_WHITE

    # Middle: 5 numbered steps
    seq_top = banner_h + Inches(0.25)
    seq_h = Inches(4.4)
    steps = [
        ('1', 'Share your framework. Ask what is missing.'),
        ('2', 'Decide what data your framework needs. Pull what you can from the one-pager. Ask AI for what is missing.'),
        ('3', 'Pick your biggest uncertainty. Ask what the call said about it.'),
        ('4', 'Ask it to argue the strongest case against your view.'),
        ('5', 'Land it. Buy / Hold / Sell with one-sentence conviction and uncertainty.'),
    ]
    step_h = Inches(0.78)
    step_gap = Inches(0.08)
    step_top = seq_top
    for num, text in steps:
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), step_top + Inches(0.1), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = AG_BLUE
        badge.line.fill.background()
        btf = badge.text_frame
        btf.margin_top = Inches(0)
        btf.margin_bottom = Inches(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = num
        br.font.name = 'Montserrat'
        br.font.size = Pt(20)
        br.font.bold = True
        br.font.color.rgb = AG_WHITE
        # Step text
        add_rich_text(
            slide,
            Inches(1.45), step_top + Inches(0.12),
            Inches(11.4), step_h,
            [{
                'text': text,
                'size': 17, 'bold': True, 'color': AG_DARK,
            }],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        step_top += step_h + step_gap

    # Bottom strip: form-a-view + reassurance + QR
    bottom_top = Inches(6.05)
    # Left: form-a-view + reassurance
    add_rich_text(
        slide,
        Inches(0.7), bottom_top,
        Inches(10.0), Inches(0.95),
        [
            {'text': 'Deliverable: Buy / Hold / Sell + one sentence conviction + one sentence uncertainty.',
             'size': 12, 'bold': True, 'color': AG_DARK, 'space_after': 3},
            {'text': 'You will not finish all five. That is fine. The five-step structure is what you take home. Pick it up at any step and your AI session improves.',
             'size': 10, 'italic': True, 'color': AG_DARK_GRAY},
        ]
    )
    # Right: QR (small, positioned to avoid page number)
    add_qr_placeholder(
        slide,
        Inches(11.7), Inches(5.7),
        Inches(0.85),
        'PANW one-pager'
    )


def build_slide_8_pair_post(slide):
    """Slide 8: Pair + Post standing. Two steps + QR."""
    strip_slide_keep_footer(slide)
    add_blue_header_bar(slide, 'Pair, then post', 'Compare. Land one view. Drop it in the feed.')

    # Left: two numbered steps
    steps = [
        ('1', 'Compare with your neighbor.', 'Land on one combined view. Disagreement is fine. Pick a deliverable.'),
        ('2', 'Post the combined view.', 'Buy / Hold / Sell. One sentence conviction. One sentence uncertainty.'),
    ]
    step_top = Inches(1.5)
    step_h = Inches(2.0)
    for num, head, body in steps:
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), step_top + Inches(0.2), Inches(0.7), Inches(0.7))
        badge.fill.solid()
        badge.fill.fore_color.rgb = AG_BLUE
        badge.line.fill.background()
        btf = badge.text_frame
        btf.margin_top = Inches(0)
        btf.margin_bottom = Inches(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = num
        br.font.name = 'Montserrat'
        br.font.size = Pt(26)
        br.font.bold = True
        br.font.color.rgb = AG_WHITE
        # Text
        add_rich_text(
            slide,
            Inches(1.6), step_top + Inches(0.2),
            Inches(7.5), step_h,
            [
                {'text': head, 'size': 20, 'bold': True, 'color': AG_DARK, 'space_after': 4},
                {'text': body, 'size': 14, 'color': AG_DARK_GRAY},
            ]
        )
        step_top += step_h + Inches(0.2)

    # Right: QR placeholder
    add_qr_placeholder(
        slide,
        Inches(9.6), Inches(2.0),
        Inches(3.0),
        'Mentimeter Beat 3 posting board'
    )

    # Bottom callout
    callout = add_rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5), fill=AG_LIGHT_BLUE)
    tf = callout.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Anonymous. No name. No login.'
    r.font.name = 'Montserrat'
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = AG_DARK


def build_slide_9_debrief(slide):
    """Slide 9: Debrief. Three blocks: Feed vs demo, Jagged edge, Q3 revisited."""
    strip_slide_keep_footer(slide)
    add_blue_header_bar(slide, 'Debrief', 'Where the room landed. Where the AI-assisted view landed. Where the edge is.')

    # Three columns
    col_top = Inches(1.4)
    col_h = Inches(4.6)
    margin = Inches(0.4)
    gap = Inches(0.3)
    total_w = SLIDE_W - margin * 2
    col_w = Emu(int((total_w - gap * 2) / 3))

    blocks = [
        {
            'badge': '01',
            'title': 'Feed vs demo',
            'prompts': [
                'Where did the room agree with the AI-assisted view?',
                'Where did it disagree?',
                'What surprised you?',
            ],
        },
        {
            'badge': '02',
            'title': 'The jagged edge',
            'prompts': [
                'Where did AI add real value?',
                'Where was human judgment load-bearing?',
                'What would you do differently next time?',
            ],
        },
        {
            'badge': '03',
            'title': 'Q3 revisited',
            'prompts': [
                'Did the room shift on Q3 from the settling poll?',
                'If yes: what moved you?',
                'If no: why are you holding your prior?',
            ],
            'qr': True,
        },
    ]

    x = margin
    for blk in blocks:
        # Header
        badge_h = Inches(0.5)
        badge = add_rect(slide, x, col_top, col_w, badge_h, fill=AG_BLUE)
        btf = badge.text_frame
        btf.margin_left = Inches(0.2)
        btf.margin_top = Inches(0)
        btf.margin_bottom = Inches(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run()
        br.text = f"{blk['badge']}  ·  {blk['title']}"
        br.font.name = 'Montserrat'
        br.font.size = Pt(15)
        br.font.bold = True
        br.font.color.rgb = AG_WHITE
        # Body
        body_top = col_top + badge_h + Inches(0.05)
        body_h = col_h - badge_h - Inches(0.05)
        body = add_rect(slide, x, body_top, col_w, body_h, fill=AG_OFFWHITE)
        tf = body.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, prompt in enumerate(blk['prompts']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(10)
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = prompt
            r.font.name = 'Montserrat'
            r.font.size = Pt(13)
            r.font.color.rgb = AG_DARK
        # QR if applicable (for Q3)
        if blk.get('qr'):
            # Add QR placeholder in the bottom of this column (well inside box)
            qr_size = Inches(1.1)
            qr_left = x + col_w - qr_size - Inches(0.2)
            qr_top = body_top + body_h - qr_size - Inches(0.5)
            add_qr_placeholder(
                slide,
                qr_left, qr_top,
                qr_size,
                'Mentimeter Q3 re-poll',
            )

        x = x + col_w + gap

    # Bottom callout
    add_rich_text(
        slide,
        Inches(0.4), Inches(6.4),
        Inches(12.5), Inches(0.5),
        [{
            'text': 'Notice whether you led with your framework or let AI set the agenda. That is the move that compounds.',
            'size': 13, 'italic': True, 'bold': True, 'color': AG_DARK, 'align': PP_ALIGN.CENTER,
        }]
    )


def build_slide_10_take_home(slide):
    """Slide 10: Take-home CTA. Text-dominant, single line."""
    strip_slide_keep_footer(slide)
    # No blue bar; minimalist treatment
    add_rich_text(
        slide,
        Inches(0.7), Inches(0.5),
        Inches(12.0), Inches(0.7),
        [{
            'text': 'Take this home',
            'size': 22, 'bold': True, 'color': AG_BLUE,
        }]
    )

    # The CTA in the middle, large, well-spaced
    add_rich_text(
        slide,
        Inches(1.0), Inches(2.4),
        Inches(11.3), Inches(2.5),
        [
            {'text': 'Next time you use AI on a real problem,',
             'size': 26, 'color': AG_DARK_GRAY, 'align': PP_ALIGN.CENTER, 'space_after': 6},
            {'text': 'notice whether you led with your framework',
             'size': 32, 'bold': True, 'color': AG_DARK, 'align': PP_ALIGN.CENTER, 'space_after': 4},
            {'text': 'or let AI set the agenda.',
             'size': 32, 'bold': True, 'color': AG_BLUE, 'align': PP_ALIGN.CENTER},
        ]
    )

    # Bottom: a quiet note
    add_rich_text(
        slide,
        Inches(0.7), Inches(6.0),
        Inches(11.9), Inches(0.5),
        [{
            'text': 'Thank you. We are on the next slide.',
            'size': 14, 'italic': True, 'color': AG_DARK_GRAY, 'align': PP_ALIGN.CENTER,
        }]
    )


def update_slide_11_contact(slide):
    """Slide 11: Contact. Kept mostly static, just update page number."""
    # Don't modify content; AG guideline says keep static
    pass


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    p = Presentation(INPUT)
    assert len(p.slides) == 11, f'Expected 11 slides, got {len(p.slides)}'

    # Slide 1: already edited via XML
    # Slide 2: Settling poll
    build_slide_2_settling_poll(p.slides[1])
    # Slide 3: 5 components
    build_slide_3_five_components(p.slides[2])
    # Slide 4: Call and question
    build_slide_4_call_and_question(p.slides[3])
    # Slide 5: Three Beats
    build_slide_5_three_beats(p.slides[4])
    # Slide 6: Beat 1
    build_slide_6_beat_1(p.slides[5])
    # Slide 7: Beat 2
    build_slide_7_beat_2(p.slides[6])
    # Slide 8: Pair + Post
    build_slide_8_pair_post(p.slides[7])
    # Slide 9: Debrief
    build_slide_9_debrief(p.slides[8])
    # Slide 10: Take-home
    build_slide_10_take_home(p.slides[9])
    # Slide 11: Contact (no changes)

    # Update page numbers across slides 2-11
    for i, slide in enumerate(p.slides, 1):
        if i == 1:
            continue  # Title has no page number
        find_and_clear_page_number_text(slide, i)
        update_footer_page(slide, i)

    p.save(OUTPUT)
    print(f'Saved {OUTPUT}')


if __name__ == '__main__':
    main()
